import streamlit as st
import pandas as pd
import csv
import io
import plotly.express as px
import plotly.graph_objects as go
import numpy as np
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import re

# --- Page Configuration ---
st.set_page_config(page_title="converterPRO", page_icon="💡", layout="wide")

# --- Custom Professional Color Palette ---
LAB_COLORS = ['#0b41cd', '#009688', '#ff9800', '#673ab7', '#e91e63', '#00bcd4', '#4caf50', '#ffc107', '#3f51b5', '#795548', '#607d8b', '#f44336']

# --- Comprehensive Clinical Knowledge Base (Roche Pure & Pro) ---
ALARM_MAP = {
    ">Abs": {"name": "ABS over", "sev": "High", "msg": "Detected foam or air aspiration, or absorbance value exceeded 3.3. Check for sample integrity.", "type": "Analytical", "action": "Check sample for bubbles/foam. Review reaction curve."},
    "ADC.E": {"name": "ADC abnormal", "sev": "Critical", "msg": "The ADC value of the primary or secondary wavelength is zero, or ISE cannot read ADC data.", "type": "Hardware", "action": "Hardware check required. Contact service if persistent."},
    ">Cuvet": {"name": "ABS Cell blank abnormal", "sev": "Medium", "msg": "Cell blank value used for measurement deviates by more than 0.1Abs.", "type": "Analytical", "action": "Perform cell wash or replace cuvette."},
    ">Reac0": {"name": "Reaction limit over", "sev": "High", "msg": "Endpoint assay measuring point out of defined reaction limit, or all rate assay points exceed limit.", "type": "Analytical", "action": "Review sample dilution or potential interference."},
    ">Reac1": {"name": "Reaction limit over (2nd..)", "sev": "High", "msg": "The second and subsequent points exceed the reaction limit.", "type": "Analytical", "action": "Review reaction curve."},
    ">Reac2": {"name": "Reaction limit over (3rd..)", "sev": "High", "msg": "The third and subsequent points exceed the reaction limit.", "type": "Analytical", "action": "Review reaction curve."},
    ">Lin": {"name": "Linearity abnormal 1/2", "sev": "High", "msg": "In rate assay, the reaction linearity exceeds the specified limit value.", "type": "Analytical", "action": "Review sample dilution requirement."},
    ">Kin": {"name": "Kinetic unstable", "sev": "Medium", "msg": "Detected by Kinetic unstable check.", "type": "Analytical", "action": "Review reaction kinetics."},
    ">Kin1": {"name": "Kinetic unstable 1", "sev": "Medium", "msg": "Detected by Kinetic unstable check.", "type": "Analytical", "action": "Review reaction kinetics."},
    ">Kin2": {"name": "Kinetic unstable 2", "sev": "Medium", "msg": "Detected by Kinetic unstable check.", "type": "Analytical", "action": "Review reaction kinetics."},
    ">Kin3": {"name": "Kinetic unstable 3", "sev": "Medium", "msg": "Detected by Kinetic unstable check.", "type": "Analytical", "action": "Review reaction kinetics."},
    "Samp.?": {"name": "ABS maximum over (non-linear curve)", "sev": "High", "msg": "Sample absorbance is outside the theoretical maximum/minimum of the calibration curve.", "type": "Analytical", "action": "Dilute and rerun."},
    "Rough": {"name": "Kinetic Roughness Check", "sev": "Medium", "msg": "Detected by Kinetic Roughness Check.", "type": "Analytical", "action": "Review assay parameters."},
    "Hook": {"name": "High Dose Hook Effect Check", "sev": "High", "msg": "Antigen excess detected.", "type": "Analytical", "action": "Dilute sample and rerun."},
    "Samp.S": {"name": "Sample short", "sev": "Critical", "msg": "Liquid level not detected before aspiration. Check sample volume.", "type": "Pre-Analytical", "action": "Add more sample or check for micro-cups."},
    "Samp.C": {"name": "Sample clot", "sev": "Critical", "msg": "Specified volume not aspirated due to clogging or short sample.", "type": "Pre-Analytical", "action": "Re-spin sample to remove fibrin/clots."},
    "Samp.B": {"name": "Sample air bubble", "sev": "High", "msg": "Air bubble detected during aspiration.", "type": "Pre-Analytical", "action": "Remove bubbles from sample surface."},
    "Samp.O": {"name": "Sample carry over", "sev": "Low", "msg": "Sample carry over wash required.", "type": "Hardware", "action": "System will automatically wash. Monitor for patterns."},
    "Samp.V": {"name": "Sample height abnormal", "sev": "High", "msg": "Liquid level detected over 57mm from bottom. Tube may be overfilled.", "type": "Pre-Analytical", "action": "Check primary tube volume/height."},
    "SASP.A": {"name": "Sample probe pressure abnormal", "sev": "Critical", "msg": "Pressure abnormality detected in sample probe flow path.", "type": "Hardware", "action": "Check probe for clogs/bends. Perform probe wash."},
    "S2PL.E": {"name": "Sample probe (S2) pressure abnormal", "sev": "Critical", "msg": "Pressure abnormality detected after S2 probe pipetting.", "type": "Hardware", "action": "Check S2 probe for clogs/bends."},
    "SLLD.E": {"name": "Sample LLD abnormal", "sev": "High", "msg": "LLD failed to start or complete. Check for dirt on tip.", "type": "Hardware", "action": "Clean probe tip. Check for deep clots."},
    "SLLD.N": {"name": "Sample LLD noise", "sev": "Medium", "msg": "Liquid surface detection failed due to bubbles or static electricity.", "type": "Pre-Analytical", "action": "Remove bubbles. Check tube material."},
    "Det.S": {"name": "Carry over detergent short", "sev": "High", "msg": "Shortage of detergent for reagent carryover evasion.", "type": "Hardware", "action": "Replenish system detergent."},
    "CarOvr": {"name": "Potential carry over", "sev": "Medium", "msg": "Signal level is abnormally low, suggesting potential carryover.", "type": "Analytical", "action": "Review preceding samples on track."},
    "Reag.S": {"name": "Reagent short", "sev": "Critical", "msg": "Liquid level not detected in reagent container. Foam/air detected.", "type": "Reagent", "action": "Replace reagent pack."},
    "ReagEx": {"name": "Reagent Expired Date", "sev": "High", "msg": "Expired reagent used; result not guaranteed.", "type": "Reagent", "action": "Replace with unexpired reagent."},
    "Reag.H": {"name": "Reagent hovering", "sev": "Medium", "msg": "Probe hovers over reaction disk.", "type": "Hardware", "action": "Mechanical check needed."},
    "Reag.F": {"name": "Reagent film detection", "sev": "Medium", "msg": "Probe detects film on reagent/dilution/pretreatment/ProCell/CleanCell.", "type": "Reagent", "action": "Check reagent surface. Discard if contaminated."},
    "Reag.T": {"name": "Reagent disk temperature", "sev": "Critical", "msg": "Reagent disk temperature out of range.", "type": "Hardware", "action": "Check lab ambient temp. Call service."},
    "Inc.T": {"name": "Incubator temperature", "sev": "Critical", "msg": "Incubator temperature out of range.", "type": "Hardware", "action": "Call service immediately."},
    "SysR.T": {"name": "System reagent temperature", "sev": "High", "msg": "ProCell/CleanCell temperature out of range.", "type": "Hardware", "action": "Check system fluid temperatures."},
    "Cell.T": {"name": "Cell temperature", "sev": "High", "msg": "Measuring cell temperature out of range.", "type": "Hardware", "action": "Call service."},
    "WBSS.T": {"name": "Washing buffer SS temperature", "sev": "Medium", "msg": "PreClean separation station temperature out of range.", "type": "Hardware", "action": "Monitor system."},
    "WB.T": {"name": "Washing buffer temperature", "sev": "Medium", "msg": "PreClean temperature out of range.", "type": "Hardware", "action": "Monitor system."},
    "ISE.N": {"name": "ISE noise error", "sev": "High", "msg": "Fluctuation in electromotive force exceeds limits (Na: 0.7mV, K: 1.0mV, Cl: 0.8mV).", "type": "Analytical", "action": "Perform ISE prime. Check electrodes."},
    "ISE.E": {"name": "ISE voltage level error", "sev": "High", "msg": "Mean EMF of internal reference out of range.", "type": "Analytical", "action": "Check internal standard fluid."},
    "ElecEx": {"name": "Expired ISE electrode", "sev": "High", "msg": "Expired electrode used; result not guaranteed.", "type": "Hardware", "action": "Replace ISE electrodes."},
    "OBS.EL": {"name": "On board stability/count of ISE exceeded", "sev": "Medium", "msg": "On board stability time or count exceeded for ISE electrodes.", "type": "Hardware", "action": "Replace ISE electrodes."},
    "OBS.RR": {"name": "On board stability limit over on reagents", "sev": "Medium", "msg": "OBS limit exceeded on Reagent Rotor and ISE reagents.", "type": "Reagent", "action": "Replace aged reagents."},
    ">Test": {"name": "Technical Limit over (upper)", "sev": "High", "msg": "Concentration exceeds technical limit/measuring range.", "type": "Analytical", "action": "Dilute sample and rerun."},
    "<Test": {"name": "Technical Limit over (lower)", "sev": "High", "msg": "Concentration is below technical limit/measuring range.", "type": "Analytical", "action": "Report as < Technical Limit."},
    ">Rept": {"name": "Repeat limit over (upper)", "sev": "Medium", "msg": "Result exceeds upper limit of specified repeat range.", "type": "Analytical", "action": "Review result. May auto-repeat."},
    "<Rept": {"name": "Repeat limit over (lower)", "sev": "Medium", "msg": "Result falls below lower limit of specified repeat range.", "type": "Analytical", "action": "Review result."},
    "H": {"name": "Above expected value", "sev": "Low", "msg": "Result is higher than reference range.", "type": "Clinical", "action": "Clinical review."},
    "L": {"name": "Below expected value", "sev": "Low", "msg": "Result is lower than reference range.", "type": "Clinical", "action": "Clinical review."},
    "Calc.?": {"name": "Calculation not possible", "sev": "High", "msg": "Denominator is zero, overflow in exponential calc, or result left blank.", "type": "Software", "action": "Review raw data points."},
    "ClcT.E": {"name": "Calculation test error", "sev": "High", "msg": "Data alarm occurred for a test needed in a calculated result.", "type": "Software", "action": "Review base test results."},
    "Over.E": {"name": "Overflow", "sev": "High", "msg": "Output figure exceeds defined digits.", "type": "Software", "action": "Check LIS transmission settings."},
    "eflowE": {"name": "e flow error", "sev": "Medium", "msg": "Sub result measured in e flow has data alarm.", "type": "Software", "action": "Review sub-results."},
    "eflowW": {"name": "e flow warning", "sev": "Low", "msg": "Higher Uncertainty flag attached to a Sub Result.", "type": "Software", "action": "Review sub-results."},
    "HU": {"name": "Higher uncertainty", "sev": "Low", "msg": "Result is between Technical Limit Low and Higher Uncertainty Limit.", "type": "Analytical", "action": "Result is valid but close to limits."},
    "<SigL": {"name": "Low level signal", "sev": "High", "msg": "Effective signal lower than specified limit.", "type": "Analytical", "action": "Check reagent/calibration."},
    ">Curr": {"name": "Current range over", "sev": "High", "msg": "Measuring cell current out of range in determination cycle.", "type": "Hardware", "action": "Hardware check."},
    "QCErr": {"name": "QC error", "sev": "High", "msg": "Error related to QC measurement.", "type": "QC", "action": "Review QC rules and calibration."},
    "QCLow": {"name": "QC out of range (Low)", "sev": "High", "msg": "QC result violates low limit rules.", "type": "QC", "action": "Review QC rules and calibration."},
    "QCHigh": {"name": "QC out of range (High)", "sev": "High", "msg": "QC result violates high limit rules.", "type": "QC", "action": "Review QC rules and calibration."},
    ">I.L": {"name": "Lipemia index interference", "sev": "Medium", "msg": "Lipemia value exceeds specified limit.", "type": "Pre-Analytical", "action": "Ultracentrifuge sample and rerun."},
    ">I.H": {"name": "Hemolysis index interference", "sev": "Medium", "msg": "Hemolysis value exceeds specified limit.", "type": "Pre-Analytical", "action": "Request new sample redraw."},
    ">I.I": {"name": "Icteric index interference", "sev": "Medium", "msg": "Icteric value exceeds specified limit.", "type": "Pre-Analytical", "action": "Note clinical condition."},
    ">I.LH": {"name": "Lipemia/Hemolysis interference", "sev": "High", "msg": "Both lipemia and hemolysis exceed limits.", "type": "Pre-Analytical", "action": "Request redraw. Difficult to salvage."},
    ">I.LI": {"name": "Lipemia/Icteric interference", "sev": "High", "msg": "Both lipemia and icteric exceed limits.", "type": "Pre-Analytical", "action": "Ultracentrifuge and note icterus."},
    ">I.HI": {"name": "Hemolysis/Icteric interference", "sev": "High", "msg": "Both hemolysis and icteric exceed limits.", "type": "Pre-Analytical", "action": "Request redraw."},
    ">I.LHI": {"name": "Lipemia/Hemolysis/Icteric interference", "sev": "Critical", "msg": "All HIL values exceed limits.", "type": "Pre-Analytical", "action": "Sample heavily compromised. Redraw."},
    "na.LHI": {"name": "Sample Index not performed", "sev": "Low", "msg": "Sample index measurement could not be performed.", "type": "Analytical", "action": "Ensure HIL mapping is active."},
    "Cal.E": {"name": "CALIB error", "sev": "High", "msg": "Calibrator concentration differs from limits or previous calibration failed.", "type": "Calibration", "action": "Review calibrator prep. Recalibrate."},
    "Call": {"name": "Calibration result invalid", "sev": "High", "msg": "Result generated with invalid transferred calibration.", "type": "Calibration", "action": "Recalibrate immediately."},
    "Dup.E": {"name": "DUPLICATE error", "sev": "High", "msg": "Difference between 1st and 2nd cal measurement is out of range.", "type": "Calibration", "action": "Check calibrator homogeneity."},
    "Std.E": {"name": "STANDARD error", "sev": "Critical", "msg": "Calibration failed due to severe hardware/fluidic alarm (clot, bubble, etc.).", "type": "Calibration", "action": "Resolve underlying hardware error. Recalibrate."},
    "Sens.E": {"name": "SENSITIVITY error", "sev": "High", "msg": "Sensitivity check failed for linear/nonlinear calibration.", "type": "Calibration", "action": "Check reagent lot and calibrator concentrations."},
    "SD.E": {"name": "SD limit error", "sev": "High", "msg": "SD value larger than specified limit during calibration.", "type": "Calibration", "action": "Check calibrator stability/probe precision."},
    "Slop.E": {"name": "Slope abnormal", "sev": "High", "msg": "ISE slope value out of range.", "type": "Calibration", "action": "Perform ISE maintenance."},
    "IStd.E": {"name": "IS concentration abnormal", "sev": "High", "msg": "Internal Standard concentration out of range.", "type": "Calibration", "action": "Replace Internal Standard fluid."},
    "Rsp1.E": {"name": "Response(ISE) abnormal 1", "sev": "High", "msg": "A Factor outside limit.", "type": "Calibration", "action": "Perform ISE prime."},
    "Rsp2.E": {"name": "Response(ISE) abnormal 2", "sev": "High", "msg": "A Factor outside limit.", "type": "Calibration", "action": "Clean ISE flow path."},
    "S1A.E": {"name": "S1ABS abnormal", "sev": "High", "msg": "Expected absorbance outside S1 Abs Limit during calibration.", "type": "Calibration", "action": "Check blank/Std 1 integrity."},
    "Diff.E": {"name": "Minimum acceptable difference", "sev": "High", "msg": "Signal difference between calibrator levels is below permissible value.", "type": "Calibration", "action": "Check reagent viability."}
}

@st.cache_data
def process_data(file_bytes):
    content = file_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(content))
    raw_data = list(reader)
    if len(raw_data) < 2: return None
    row0, header_row = raw_data[0], raw_data[1]

    try:
        draw_idx = header_row.index("Drawing_Date_Time")
        cup_idx = header_row.index("Sample_Cup")
        start_col = header_row.index("Result") 
        block_size = 11 if "EMF1" in header_row else 8
        
        block_template = header_row[start_col : start_col + block_size]
        module_sub_idx = next((i for i, h in enumerate(block_template) if h in ["Module", "AU", "Unit_ID"]), -1)
        alarm_sub_idx = next((i for i, h in enumerate(block_template) if h in ["Data_Alarm", "Alarm"]), -1)
        sampling_sub_idx = next((i for i, h in enumerate(block_template) if "Sampling_Date_Time" in h), -1)

        fixed_before = header_row[:draw_idx + 1]
        comment_cols = header_row[draw_idx + 1 : cup_idx]
        fixed_after = header_row[cup_idx : start_col]
        
        standard_block_headers = list(block_template)
        if module_sub_idx != -1: standard_block_headers[module_sub_idx] = "Module"
        if alarm_sub_idx != -1: standard_block_headers[alarm_sub_idx] = "Data_Alarm"
        if sampling_sub_idx != -1: standard_block_headers[sampling_sub_idx] = "Sampling_Date_Time"
        
        final_headers = fixed_before + comment_cols + fixed_after + ["ACN code", "Parameter"] + standard_block_headers
        
        blocks = []
        for col in range(start_col, len(header_row), block_size):
            if col + block_size <= len(header_row):
                acn, param = (row0[col], row0[col+1]) if col+1 < len(row0) else ("", "")
                for row_vals in raw_data[2:]:
                    if len(row_vals) > col and any(val.strip() for val in row_vals[col:col+block_size]):
                        payload = row_vals[:draw_idx+1] + row_vals[draw_idx+1:cup_idx] + row_vals[cup_idx:start_col] + [acn, param] + row_vals[col:col+block_size]
                        if len(payload) == len(final_headers): blocks.append(payload)
        
        df = pd.DataFrame(blocks, columns=final_headers)
        df['Arrived_Date_Time'] = pd.to_datetime(df['Arrived_Date_Time'], errors='coerce')
        df['Sampling_Date_Time'] = pd.to_datetime(df['Sampling_Date_Time'], errors='coerce')
        df['Result_Numeric'] = pd.to_numeric(df['Result'], errors='coerce')

        def parse_au(au_str):
            if not au_str or pd.isna(au_str): return "N/A", "Unknown", "Unknown"
            au_str = str(au_str).strip()
            if '-' in au_str:
                parts = au_str.split('-')
                pos, mtype = parts[0], (parts[1] if len(parts) > 1 else "N/A")
                sub = parts[2] if len(parts) > 2 else "0"
                return pos, mtype, f"{mtype}-{sub}"
            else:
                return "1", au_str, f"{au_str}-0"

        df[['AU_Pos', 'AU_Class', 'AU_SubUnit']] = df['Module'].apply(lambda x: pd.Series(parse_au(x)))
        
        def map_alarm_type(c):
            if pd.isna(c): return "None"
            c_str = str(c).strip()
            
            # HARD FIX: Ignore standalone raw numbers and treat them as 'None' (No error)
            if not c_str or c_str.lstrip('-').replace('.','',1).isdigit(): 
                return "None"
                
            if c_str in ALARM_MAP: return ALARM_MAP[c_str]["type"]
            
            c_up = c_str.upper()
            if "QC" in c_up: return "QC"
            if "SAMP" in c_up: return "Pre-Analytical"
            if "REAG" in c_up or "OBS" in c_up: return "Reagent"
            if "CAL" in c_up or "STD" in c_up or "DUP" in c_up: return "Calibration"
            if "ISE" in c_up or "ABS" in c_up or "KIN" in c_up or "REAC" in c_up: return "Analytical"
            return "Unknown"

        df['Alarm_Code'] = df['Data_Alarm'].apply(lambda x: "" if pd.isna(x) else str(x).strip())
        df['Alarm_Type'] = df['Alarm_Code'].apply(map_alarm_type)
        df['Alarm_Meaning'] = df['Alarm_Code'].apply(lambda x: ALARM_MAP.get(x, {"name": x})['name'] if x else "")

        mappings = {
            "Gender": {"0": "Not entered", "1": "Male", "2": "Female"},
            "Discrimination": {"1": "Patient (Routine)", "2": "Patient (STAT)", "3": "QC (Control)"},
            "Run": {"1": "1st run", "2": "Rerun"}
        }
        for col in mappings:
            if col in df.columns: df[col] = df[col].astype(str).map(mappings[col]).fillna(df[col])
            
        return df
    except Exception as e:
        st.error(f"Processing Error: {e}")
        return None

def render_insight(title, obs, impact, pre, status="info"):
    colors = {"info": "#0b41cd", "warning": "#ff9800", "critical": "#f44336", "success": "#4caf50"}
    st.markdown(f"""<div style="background-color: #ffffff; padding: 20px; border-radius: 10px; border-left: 10px solid {colors.get(status)}; margin-top: 10px; margin-bottom: 30px; border: 1px solid #eee; box-shadow: 2px 2px 5px rgba(0,0,0,0.05);">
        <h4 style="margin-top:0; color: {colors.get(status)};">🔎 Insight Card: {title}</h4>
        <div style="display: flex; gap: 20px;"><div style="flex: 1;"><strong>Observation</strong><br>{obs}</div>
        <div style="flex: 1;"><strong>Impact</strong><br>{impact}</div><div style="flex: 1;"><strong>Action/Checklist</strong><br>{pre}</div></div></div>""", unsafe_allow_html=True)

def create_ppt(figs_dict):
    prs = Presentation()
    for title, fig in figs_dict.items():
        fig.layout.paper_bgcolor = '#ffffff'
        fig.layout.plot_bgcolor = '#ffffff'
        fig.update_layout(template="plotly_white", font=dict(color="#000000"))
        
        img_bytes = fig.to_image(format="png", engine="kaleido", width=1000, height=550, scale=2)
        img_stream = io.BytesIO(img_bytes)
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.5), width=Inches(9))
        
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

# --- Sidebar ---
with st.sidebar:
    st.title("💡 converterPRO")
    uploaded_file = st.file_uploader("Upload Instrument CSV", type=["csv"])
    if uploaded_file:
        raw_df = process_data(uploaded_file.getvalue())
        if raw_df is not None:
            st.markdown("---")
            st.subheader("📅 Filter View")
            min_d, max_d = raw_df['Arrived_Date_Time'].min().date(), raw_df['Arrived_Date_Time'].max().date()
            sel_range = st.date_input("Date Range", [min_d, max_d], min_value=min_d, max_value=max_d)
            sel_cats = st.multiselect("Data Categories", raw_df['Discrimination'].unique().tolist(), default=raw_df['Discrimination'].unique().tolist())
            
            st.markdown("---")
            st.subheader("🛡️ Privacy & Compliance")
            scrub_phi = st.checkbox("Scrub Patient Data (PHI)", value=True, help="Replaces text in comment fields with [REDACTED] to protect patient privacy.")
    
    st.markdown("---")
    st.caption("⚙️ **Engine Details**")
    st.markdown("- **Adaptive Engine:** v9.3\n- **Compatibility:** cobas pro / cobas pure\n- **Status:** Validated")
    st.markdown("---")
    st.markdown("© 2026 **LabMesh.com**")

export_figs = {}

# --- Main App ---
if uploaded_file and 'raw_df' in locals() and raw_df is not None:
    start_d, end_d = sel_range[0], (sel_range[1] if len(sel_range) > 1 else sel_range[0])
    mask = (raw_df['Arrived_Date_Time'].dt.date >= start_d) & (raw_df['Arrived_Date_Time'].dt.date <= end_d) & (raw_df['Discrimination'].isin(sel_cats))
    
    df = raw_df.loc[mask].copy()
    if scrub_phi:
        for col in df.columns:
            if 'Comment' in col:
                df[col] = df[col].apply(lambda x: "[REDACTED]" if pd.notna(x) and str(x).strip() != "" else x)

    t = st.tabs(["📄 Raw Data", "📊 Throughput", "🧪 Quality Control", "🔄 Reruns", "⚠️ Alarms & Risk", "⚙️ Hardware Load"])
    
    with t[0]:
        st.dataframe(df, use_container_width=True)
        st.download_button("📥 Export CSV", df.to_csv(index=False), f"Enriched_{uploaded_file.name}")

    with t[1]:
        st.subheader("Laboratory Throughput & Peak Stress")
        # Filter out 'Unknown' because they are software calculations, not mechanical loads
        u_df = df.dropna(subset=['Sampling_Date_Time']).copy()
        u_df = u_df[u_df['AU_Class'] != "Unknown"]
        
        if not u_df.empty:
            u_df['S_Hour'] = u_df['Sampling_Date_Time'].dt.hour
            u_df['S_Date'] = u_df['Sampling_Date_Time'].dt.date.astype(str) 
            h_util = u_df.groupby(['S_Date', 'S_Hour', 'AU_Class']).size().reset_index(name='Tests')
            
            sel_m = st.selectbox("View sampling pattern for:", h_util['AU_Class'].unique().tolist())
            fig_line = px.line(h_util[h_util['AU_Class'] == sel_m], x='S_Hour', y='Tests', color='S_Date', text='Tests', markers=True, color_discrete_sequence=LAB_COLORS, title=f"Instrument Sampling Rate (Tests/Hr): {sel_m}")
            fig_line.update_traces(textposition="top center")
            fig_line.update_layout(xaxis=dict(tickmode='linear', tick0=0, dtick=1, range=[0, 23]))
            st.plotly_chart(fig_line, use_container_width=True)
            export_figs["Instrument Sampling Rate"] = fig_line
            
            is_pure = any("303" in str(x) or "402" in str(x) for x in h_util['AU_Class'])
            
            if is_pure:
                caps = {"c 303": (450, 360), "ISE": (450, 360), "e 402": (120, 100), "EFLW": (0,0)}
            else:
                caps = {"c 503": (1000, 800), "c 703": (2000, 1800), "ISE": (900, 850), "e 801": (300, 275), "EFLW": (0,0)}
            
            peak_stats = []
            for m_type, (m_max, m_prac) in caps.items():
                if any(m_type in str(x) for x in h_util['AU_Class']) and m_max > 0:
                    peak_val = h_util[h_util['AU_Class'].str.contains(m_type, na=False)]['Tests'].max()
                    peak_stats.append({'Module': m_type, 'Peak': peak_val, 'Prac': m_prac, 'Theo': m_max})
            
            if peak_stats:
                fig_p = go.Figure()
                fig_p.add_trace(go.Bar(x=[d['Module'] for d in peak_stats], y=[d['Peak'] for d in peak_stats], text=[int(d['Peak']) for d in peak_stats], textposition='auto', name="Actual Peak", marker_color='#0b41cd'))
                for i, d in enumerate(peak_stats):
                    fig_p.add_shape(type="line", x0=i-0.3, y0=d['Prac'], x1=i+0.3, y1=d['Prac'], line=dict(color="orange", width=3, dash="dash"), name="Practical Limit")
                    fig_p.add_shape(type="line", x0=i-0.3, y0=d['Theo'], x1=i+0.3, y1=d['Theo'], line=dict(color="red", width=3), name="Theoretical Limit")
                fig_p.update_layout(title="Peak Stress vs Module Capacity (Orange = Practical Limit)")
                st.plotly_chart(fig_p, use_container_width=True)
                export_figs["Peak Stress vs Capacity"] = fig_p

                stress_mod = [d['Module'] for d in peak_stats if d['Peak'] > d['Prac']]
                if stress_mod:
                    render_insight("Peak Capacity Stress", f"{', '.join(stress_mod)} exceeded practical throughput limits.", "Operating above practical limits significantly delays sample pipetting.", "Flatten the peak by batching routine non-urgent samples.", "warning")
                else:
                    render_insight("Throughput Efficiency", "All modules are operating within practical capacity.", "Workflow and TAT should remain stable without bottlenecks.", "Optimal loading rate detected.", "success")
        else: st.info("No mechanical sampling data available for throughput analysis.")
        
        st.markdown("---")
        st.subheader("Total Test Arrival Pattern (24h)")
        a_df = df.dropna(subset=['Arrived_Date_Time']).copy()
        if not a_df.empty:
            a_df['A_Hour'] = a_df['Arrived_Date_Time'].dt.hour
            a_df['A_Date'] = a_df['Arrived_Date_Time'].dt.date.astype(str)
            arr_counts = a_df.groupby(['A_Date', 'A_Hour']).size().reset_index(name='Total Tests')
            
            fig_arr = px.line(arr_counts, x='A_Hour', y='Total Tests', text='Total Tests', color='A_Date', markers=True, color_discrete_sequence=LAB_COLORS, title="Hourly Total Volume of Arriving Tests")
            fig_arr.update_traces(textposition="top center")
            fig_arr.update_layout(xaxis=dict(tickmode='linear', tick0=0, dtick=1, range=[0, 23]))
            st.plotly_chart(fig_arr, use_container_width=True)
            export_figs["Hourly Arrival Pattern"] = fig_arr
        else: st.info("No arrival data available for analysis.")

    with t[2]:
        st.subheader("QC Precision & Stability")
        q_df = df[df['Discrimination'].str.contains("QC", na=False)].copy()
        if not q_df.empty:
            q_df['HF'] = q_df['Arrived_Date_Time'].dt.hour + q_df['Arrived_Date_Time'].dt.minute/60
            fig_qc = px.scatter(q_df, x='HF', y='Parameter', color='Parameter', color_discrete_sequence=LAB_COLORS, title="QC Timing Matrix (24h)")
            st.plotly_chart(fig_qc, use_container_width=True)
            export_figs["QC Timing Matrix"] = fig_qc
            
            qc_stats = q_df.groupby(['Parameter', 'Sample_ID'])['Result_Numeric'].agg(Mean='mean', SD='std').reset_index()
            qc_stats['CV%'] = ((qc_stats['SD'] / qc_stats['Mean']) * 100).round(2)
            st.dataframe(qc_stats, use_container_width=True)
            
            c1, c2 = st.columns(2)
            with c1: 
                fig_chem = px.box(q_df[q_df['AU_Class'].str.contains("303|503|ISE", na=False)], x='Parameter', y='Result_Numeric', color='Parameter', color_discrete_sequence=LAB_COLORS, title="Chemistry Stability")
                st.plotly_chart(fig_chem, use_container_width=True)
                export_figs["Chemistry QC Stability"] = fig_chem
            with c2: 
                fig_ia = px.box(q_df[q_df['AU_Class'].str.contains("402|801", na=False)], x='Parameter', y='Result_Numeric', color='Parameter', color_discrete_sequence=LAB_COLORS, title="IA Stability")
                st.plotly_chart(fig_ia, use_container_width=True)
                export_figs["Immunoassay QC Stability"] = fig_ia
            
            bad_cv_df = qc_stats[qc_stats['CV%'] > 5]
            if not bad_cv_df.empty:
                bad_assays = ", ".join(bad_cv_df['Parameter'].unique())
                render_insight("QC Drift Detection", f"The following assays exceed 5% CV: **{bad_assays}**", "Indicates precision issues, reagent instability, or probe wear.", f"Recalibrate or perform probe maintenance on {bad_assays}.", "warning")
            else:
                render_insight("QC Status", "All parameters show stable CV% below 5%.", "Precision is within optimal technical limits.", "No action needed.", "success")
        else:
            st.info("No QC data found.")

    with t[3]:
        st.subheader("Rerun & Yield Analysis")
        r_counts = df['Run'].value_counts()
        if not r_counts.empty:
            fig_pie = px.pie(values=r_counts.values, names=r_counts.index, hole=0.5, color_discrete_map={'1st run': '#0b41cd', 'Rerun': '#f44336'}, title="First-Pass Yield vs Reruns")
            fig_pie.update_traces(textinfo='percent+value+label', textposition='inside')
            st.plotly_chart(fig_pie, use_container_width=True)
            export_figs["First-Pass Yield Ratio"] = fig_pie
            
            rerun_only = df[df['Run'] == 'Rerun']
            rerun_rate = (len(rerun_only) / len(df)) * 100 if len(df) > 0 else 0
            
            if not rerun_only.empty:
                rerun_df = rerun_only.groupby('Parameter').size().reset_index(name='Count').sort_values('Count', ascending=False)
                fig_r_bar = px.bar(rerun_df, x='Parameter', y='Count', text='Count', title="Top Rerun Assays", color_discrete_sequence=['#f44336'])
                fig_r_bar.update_traces(textposition='auto')
                st.plotly_chart(fig_r_bar, use_container_width=True)
                export_figs["Top Rerun Assays"] = fig_r_bar
                
                top_assay = rerun_df.iloc[0]['Parameter']
                if rerun_rate < 2.0:
                    render_insight("System Yield", f"Excellent First-Pass Yield. Rerun rate is {rerun_rate:.1f}%.", "Reagent waste and TAT delays are minimal.", "System is performing optimally.", "success")
                elif rerun_rate <= 5.0:
                    render_insight("Yield Efficiency Warning", f"Elevated Rerun Rate at {rerun_rate:.1f}%.", "Reruns are increasing reagent costs and TAT.", f"Investigate '{top_assay}' for frequent errors.", "warning")
                else:
                    render_insight("Yield Efficiency Critical", f"Severe Yield Bleed. Rerun rate is {rerun_rate:.1f}%.", "Reruns are doubling reagent costs and significantly delaying TAT.", f"Immediate audit of '{top_assay}' required.", "critical")
            else:
                render_insight("System Yield", "100% First-Pass Yield.", "Reagent waste is zero.", "System is performing optimally.", "success")
        else: st.info("No run data found.")

    with t[4]:
        st.subheader("Analytical Risk & Error Intelligence")
        
        err_df = df[df['Alarm_Type'] != "None"].copy()
        
        if not err_df.empty:
            type_counts = err_df['Alarm_Type'].value_counts().reset_index(name='Count')
            fig_type = px.pie(type_counts, values='Count', names='Alarm_Type', hole=0.4, title="Overall Lab Error Profile", color_discrete_sequence=LAB_COLORS)
            st.plotly_chart(fig_type, use_container_width=True)
            export_figs["Error Profile"] = fig_type

            err_bar = px.bar(err_df.groupby(['Module', 'Alarm_Code']).size().reset_index(name='C').sort_values('C', ascending=False).head(25), x='Alarm_Code', y='C', text='C', color='Module', color_discrete_sequence=LAB_COLORS, title="Top 25 System Alarms Triggered")
            err_bar.update_traces(textposition='auto')
            st.plotly_chart(err_bar, use_container_width=True)
            export_figs["System Alarms"] = err_bar
            
            pre_ana = err_df[err_df['Alarm_Type'] == 'Pre-Analytical']
            if not pre_ana.empty:
                top_pre = pre_ana['Alarm_Code'].value_counts().idxmax()
                render_insight("Pre-Analytical Bleed", f"Detected {len(pre_ana)} pre-analytical errors (Top: {top_pre}).", "Issues like Clots, Shorts, and HIL interferences lead to immediate probe damage or unreportable results.", "Audit centrifuge protocols and phlebotomy draw volumes.", "critical")
            
            reag_err = err_df[err_df['Alarm_Type'] == 'Reagent']
            if not reag_err.empty:
                render_insight("Reagent Management", f"Detected {len(reag_err)} Reagent/OBS flags.", "Using expired reagents or running low during peak hours halts the track.", "Review inventory and On-Board Stability (OBS) limits.", "warning")

            cal_err = err_df[err_df['Alarm_Type'] == 'Calibration']
            if not cal_err.empty:
                render_insight("Calibration Instability", f"Detected {len(cal_err)} Calibration failures.", "Failed calibrations prevent patient sample processing.", "Check calibrator lot expiry and reconstitution.", "critical")

            if pre_ana.empty and reag_err.empty and cal_err.empty:
                render_insight("System Health", "Minor analytical warnings detected.", "No critical operational halt alarms.", "Review reaction curves if >Reac alarms persist.", "info")

        else:
            render_insight("Alarm Status", "Zero flags detected.", "Results are analytically clean.", "Continue standard monitoring.", "success")

    with t[5]:
        st.subheader("Sub-Module Load Balancing")
        
        # Filter out 'Unknown' hardware units since they represent calculated math (no mechanical load)
        load_df = df.dropna(subset=['AU_Class', 'AU_SubUnit']).copy()
        load_df = load_df[load_df['AU_SubUnit'] != "Unknown"]
        
        if not load_df.empty:
            load_summary = load_df['AU_SubUnit'].value_counts().reset_index()
            load_summary.columns = ['Unit', 'Count']
            
            fig_load = px.bar(load_summary, x='Unit', y='Count', text='Count', color='Unit', color_discrete_sequence=LAB_COLORS, title="Mechanical Load per Sub-Unit")
            fig_load.update_traces(textposition='auto')
            st.plotly_chart(fig_load, use_container_width=True)
            export_figs["Sub-Module Load"] = fig_load
            
            imbalance_found = False
            imbalance_msgs = []
            
            for au_class in load_df['AU_Class'].unique():
                class_df = load_df[load_df['AU_Class'] == au_class]
                sub_load = class_df['AU_SubUnit'].value_counts()
                
                if len(sub_load) > 1:
                    imb = sub_load.max() / sub_load.min()
                    if imb > 1.25:
                        imbalance_found = True
                        imbalance_msgs.append(f"{au_class} ({imb:.1f}x)")
                        
            if imbalance_found:
                render_insight("Mechanical Wear Skew", f"Imbalance detected within identical modules: **{', '.join(imbalance_msgs)}**.", "Uneven wear accelerates part degradation and reduces module lifespan on specific analytical units.", "Re-map high-volume tests across parallel modules to balance the workload.", "warning")
            else:
                parallel_exists = any(len(load_df[load_df['AU_Class'] == c]['AU_SubUnit'].unique()) > 1 for c in load_df['AU_Class'].unique())
                if parallel_exists:
                    render_insight("Load Balance", "Workload is properly distributed among identical parallel modules.", "Even mechanical wear detected. Maximizing instrument lifespan.", "Mapping is optimal.", "success")
                else:
                    render_insight("Load Balance", "Single sub-modules detected per class.", "Natural test mix displayed.", "No parallel balancing required.", "info")
        else: st.info("No physical module load data found.")

    # --- PPT Export Section (Currently Disabled) ---
    # ... (PPT code remains commented out)
else:
    st.info("👈 Upload a CSV to begin.")
